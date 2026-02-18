# Local file ingestion helpers.
# This module extracts text from many file types and normalizes read failures
# into a single `FileReadError` type for the caller.

# Enable postponed evaluation of type annotations.
from __future__ import annotations

import codecs
import os
import re
from pathlib import Path

import ocr_api


# Custom exception raised for any read/extract failure in this module.
class FileReadError(RuntimeError):
    # No extra behavior needed; custom type is used for consistent error handling.
    pass


# Truncate long text while preserving a visible ellipsis suffix.
def _truncate(text: str, *, max_chars: int) -> str:
    # If max_chars is disabled/non-positive, return text unchanged.
    if max_chars <= 0:
        return text
    # If text already fits, return as-is.
    if len(text) <= max_chars:
        return text
    # Otherwise trim and append ellipsis.
    return text[: max_chars - 3].rstrip() + "..."


# Heuristic to detect likely binary payloads.
def _is_probably_binary(data: bytes) -> bool:
    # Empty payload is not considered binary.
    if not data:
        return False
    # Inspect only an initial sample for speed.
    sample = data[:8000]
    # NUL bytes strongly indicate binary content.
    if b"\x00" in sample:
        return True
    # ASCII printable set plus tab/newline/carriage return.
    allowed = set(range(32, 127)) | {9, 10, 13}
    # Count bytes that look printable.
    printable = sum(1 for b in sample if b in allowed)
    # Compute printable ratio.
    ratio = printable / max(1, len(sample))
    # Mark as binary when too few printable bytes are present.
    return ratio < 0.70


# Decode bytes to text using BOM-aware and fallback strategy.
def _decode_text_bytes(data: bytes) -> str:
    # Handle UTF-8 BOM.
    if data.startswith(codecs.BOM_UTF8):
        return data.decode("utf-8-sig", errors="replace")
    # Handle UTF-16 BOM (little/big endian).
    if data.startswith(codecs.BOM_UTF16_LE) or data.startswith(codecs.BOM_UTF16_BE):
        return data.decode("utf-16", errors="replace")
    # Handle UTF-32 BOM (little/big endian).
    if data.startswith(codecs.BOM_UTF32_LE) or data.startswith(codecs.BOM_UTF32_BE):
        return data.decode("utf-32", errors="replace")

    # Try UTF-8 first.
    for enc in ("utf-8", "utf-8-sig"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            # Continue to next candidate encoding.
            pass
    # Windows-friendly fallback that decodes any byte sequence.
    return data.decode("cp1252", errors="replace")


# Read raw bytes from disk with optional max-byte cap.
def _read_bytes(path: Path, *, max_bytes: int) -> bytes:
    # Open path in binary mode.
    try:
        with path.open("rb") as f:
            # If cap is set, read at most max_bytes.
            if max_bytes and max_bytes > 0:
                return f.read(max_bytes)
            # Otherwise read full file.
            return f.read()
    # Convert OS-level read errors to FileReadError.
    except OSError as e:
        raise FileReadError(f"Failed to read file: {path} ({e})") from e


# Convert basic HTML into plain text.
def _html_to_text(html: str) -> str:
    # Remove script/style/noscript blocks entirely.
    html = re.sub(r"(?is)<(script|style|noscript).*?>.*?</\\1>", " ", html)
    # Convert <br> tags to line breaks.
    html = re.sub(r"(?i)<br\\s*/?>", "\n", html)
    # Convert common block-closing tags to line breaks.
    html = re.sub(r"(?i)</(p|div|li|tr|h\\d|blockquote)>", "\n", html)
    # Remove all remaining tags.
    html = re.sub(r"(?s)<[^>]+>", " ", html)
    # Collapse repeated spaces/tabs.
    html = re.sub(r"[ \t]{2,}", " ", html)
    # Collapse overly large blank-line regions.
    html = re.sub(r"\n{3,}", "\n\n", html)
    # Return cleaned text.
    return html.strip()


# Extract raw text per PDF page.
def _read_pdf_pages(path: Path) -> list[str]:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception as e:  # pragma: no cover
        raise FileReadError(f"PDF support requires 'pypdf'. ({e})") from e

    # Open PDF file.
    try:
        reader = PdfReader(str(path))
    except Exception as e:
        raise FileReadError(f"Failed to open PDF: {path} ({e})") from e

    # Storage for page text.
    parts: list[str] = []
    # Optional environment cap for number of pages to read (0 = no cap).
    max_pages = int(os.getenv("AGENT_MAX_PDF_PAGES", "0"))
    # Iterate pages in order.
    for i, page in enumerate(reader.pages):
        # Stop if cap is reached.
        if max_pages > 0 and i >= max_pages:
            break
        # Extract page text defensively.
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            parts.append("")
    # Return one string per page (possibly empty).
    return parts


# Extract merged text from PDF pages.
def _read_pdf(path: Path) -> str:
    # Reuse per-page extraction.
    parts = _read_pdf_pages(path)
    # Join non-empty normalized pages with blank line separator.
    return "\n\n".join([p.strip() for p in parts if p and p.strip()]).strip()


# Read a PDF and keep character-span -> page mapping for downstream citation.
def read_pdf_with_page_spans(path: str | Path, *, max_chars: int | None = None) -> tuple[str, list[tuple[int, int, int]]]:
    # Normalize to Path.
    p = Path(path)
    # Extract raw text for each page.
    parts = _read_pdf_pages(p)
    # Piece buffer used to build merged text.
    pieces: list[str] = []
    # List of tuples: (start_char, end_char, page_number).
    spans: list[tuple[int, int, int]] = []
    # Current merged-text cursor.
    cursor = 0
    # Walk pages in order.
    for page_num, raw_page in enumerate(parts, start=1):
        # Normalize page text.
        text = (raw_page or "").strip()
        # Skip empty pages.
        if not text:
            continue
        # Add separator between pages when needed.
        if pieces:
            pieces.append("\n\n")
            cursor += 2
        # Record span start.
        start = cursor
        # Append page text.
        pieces.append(text)
        # Advance cursor by appended text length.
        cursor += len(text)
        # Record span end.
        end = cursor
        # Store mapping for this page range.
        spans.append((start, end, page_num))

    # Build merged text payload.
    merged = "".join(pieces)
    # Optional hard cap for merged text length.
    if max_chars is not None:
        cap = int(max_chars)
        if cap > 0 and len(merged) > cap:
            # Clip merged text to cap.
            merged = merged[:cap]
            # Clip spans so they remain valid after truncation.
            clipped: list[tuple[int, int, int]] = []
            for start, end, page_num in spans:
                if start >= cap:
                    break
                clipped_end = min(end, cap)
                if clipped_end > start:
                    clipped.append((start, clipped_end, page_num))
            spans = clipped

    # Return merged text with page-span metadata.
    return merged, spans


# Extract text from DOCX paragraphs, with fallback for mislabeled files.
def _read_docx(path: Path) -> str:
    # Fallback reader used when DOCX parser is unavailable or file is malformed.
    def _text_fallback() -> str | None:
        # Try to decode as plain text if possible.
        try:
            data = _read_bytes(path, max_bytes=int(os.getenv("AGENT_MAX_FILE_BYTES", "25000000")))
            if _is_probably_binary(data):
                return None
            text = _decode_text_bytes(data).strip()
            return text or None
        except Exception:
            return None

    try:
        from docx import Document  # type: ignore
    except Exception as e:  # pragma: no cover
        fallback = _text_fallback()
        if fallback is not None:
            return fallback
        raise FileReadError(f"DOCX support requires 'python-docx'. ({e})") from e

    # Parse DOCX file.
    try:
        doc = Document(str(path))
    except Exception as e:
        # Some files are mislabeled as .docx; attempt plain-text fallback.
        fallback = _text_fallback()
        if fallback is not None:
            return fallback
        raise FileReadError(f"Failed to open DOCX: {path} ({e})") from e

    # Collect non-empty paragraph text in order.
    paras = [p.text for p in doc.paragraphs if (p.text or "").strip()]
    # Return normalized joined text.
    return "\n".join(paras).strip()


# Extract text from PPTX slides.
def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:  # pragma: no cover
        raise FileReadError(f"PPTX support requires 'python-pptx'. ({e})") from e

    # Parse presentation file.
    try:
        pres = Presentation(str(path))
    except Exception as e:
        raise FileReadError(f"Failed to open PPTX: {path} ({e})") from e

    # Output text buffer.
    out: list[str] = []
    # Optional slide cap (0 = no cap).
    max_slides = int(os.getenv("AGENT_MAX_PPTX_SLIDES", "80"))
    # Iterate slides in order.
    for i, slide in enumerate(pres.slides):
        if max_slides > 0 and i >= max_slides:
            break
        # Add slide marker for context.
        out.append(f"[Slide {i+1}]")
        # Walk shapes and capture plain text when present.
        for shape in slide.shapes:
            text = getattr(shape, "text", "") or ""
            text = text.strip()
            if text:
                out.append(text)
        # Add blank separator between slides.
        out.append("")
    # Return normalized content.
    return "\n".join(out).strip()


# Extract text-like table content from XLSX.
def _read_xlsx(path: Path) -> str:
    try:
        import openpyxl  # type: ignore
    except Exception as e:  # pragma: no cover
        raise FileReadError(f"XLSX support requires 'openpyxl'. ({e})") from e

    # Open workbook in read-only/data-only mode.
    try:
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
    except Exception as e:
        raise FileReadError(f"Failed to open XLSX: {path} ({e})") from e

    # Limit extraction scope for safety/performance.
    max_sheets = int(os.getenv("AGENT_MAX_XLSX_SHEETS", "6"))
    max_rows = int(os.getenv("AGENT_MAX_XLSX_ROWS", "200"))
    max_cols = int(os.getenv("AGENT_MAX_XLSX_COLS", "30"))

    # Output text buffer.
    out: list[str] = []
    # Iterate worksheet names by workbook order.
    for si, name in enumerate(wb.sheetnames):
        if max_sheets > 0 and si >= max_sheets:
            break
        ws = wb[name]
        # Add worksheet marker.
        out.append(f"[Sheet] {name}")
        # Iterate rows.
        for ri, row in enumerate(ws.iter_rows(values_only=True)):
            if max_rows > 0 and ri >= max_rows:
                break
            cells = []
            # Iterate cells in row.
            for ci, val in enumerate(row):
                if max_cols > 0 and ci >= max_cols:
                    break
                if val is None:
                    cells.append("")
                else:
                    cells.append(str(val))
            # Keep rows with at least one non-empty value.
            if any(c.strip() for c in cells):
                out.append("\t".join(cells).rstrip())
        # Add blank separator between sheets.
        out.append("")
    # Return normalized content.
    return "\n".join(out).strip()


# Extract text from image using OCR.Space.
def _read_image_ocr(path: Path, *, language: str | None = None) -> str:
    # Read image bytes with service-size cap.
    data = _read_bytes(path, max_bytes=int(os.getenv("OCR_SPACE_MAX_BYTES", "8000000")))
    # Resolve OCR language preference.
    lang = (language or "").strip() or os.getenv("OCR_SPACE_LANGUAGE") or os.getenv("OCR_SPACE_LANG")
    # Call OCR service and map OCR-specific errors to FileReadError.
    try:
        return ocr_api.ocr_image_bytes(data, filename=path.name, language=lang)
    except ocr_api.OCRSpaceError as e:
        raise FileReadError(str(e)) from e


# Public entrypoint used by the agent to read local files safely.
def read_any_file(path: str, *, max_chars: int | None = None, ocr_language: str | None = None) -> str:
    """
    Best-effort text extraction from many file types.

    Supported (when deps are installed): pdf, docx, pptx, xlsx, images (OCR).
    Falls back to decoding as text for unknown extensions.
    """
    # Normalize requested path.
    p = Path(path)
    # Ensure file exists and is a regular file.
    if not p.exists() or not p.is_file():
        raise FileReadError(f"File not found: {path}")

    # Resolve output size limits.
    max_chars = int(os.getenv("AGENT_MAX_FILE_CHARS", "200000")) if max_chars is None else int(max_chars)
    max_bytes = int(os.getenv("AGENT_MAX_FILE_BYTES", "25000000"))

    # Dispatch by file extension.
    ext = p.suffix.lower()
    if ext == ".pdf":
        return _truncate(_read_pdf(p), max_chars=max_chars)
    if ext == ".docx":
        return _truncate(_read_docx(p), max_chars=max_chars)
    if ext == ".pptx":
        return _truncate(_read_pptx(p), max_chars=max_chars)
    if ext in {".xlsx", ".xlsm", ".xltx", ".xltm"}:
        return _truncate(_read_xlsx(p), max_chars=max_chars)
    if ext in {".html", ".htm"}:
        data = _read_bytes(p, max_bytes=max_bytes)
        if _is_probably_binary(data):
            raise FileReadError("HTML file looks binary/unreadable.")
        return _truncate(_html_to_text(_decode_text_bytes(data)), max_chars=max_chars)
    if ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}:
        return _truncate(_read_image_ocr(p, language=ocr_language), max_chars=max_chars)

    # Generic text fallback for unknown extension.
    data = _read_bytes(p, max_bytes=max_bytes)
    if _is_probably_binary(data):
        raise FileReadError(f"Unsupported binary file type: {ext or '(no extension)'}")
    return _truncate(_decode_text_bytes(data), max_chars=max_chars)
