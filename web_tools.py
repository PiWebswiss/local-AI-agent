# HTTPS-only web helper utilities.
# Provides search and content extraction for URLs used by the research mode.
from __future__ import annotations

import json
import io
import os
import re
import urllib.error
import urllib.parse
import urllib.request
from html.parser import HTMLParser
from typing import Any

import ocr_api


class WebToolError(RuntimeError):
    # Raised when web search/fetch fails in a normalized, user-friendly way.
    pass


# Normalize URLs and enforce HTTPS-only behavior.
def ensure_https_url(url: str) -> str:
    # Normalize whitespace and guard against None.
    url = (url or "").strip()
    # Reject empty inputs early.
    if not url:
        raise WebToolError("URL is empty.")

    # Parse user input once to inspect scheme/host/path.
    parsed = urllib.parse.urlparse(url)
    # If no scheme is given, default to HTTPS.
    if not parsed.scheme:
        url = "https://" + url.lstrip("/")
        parsed = urllib.parse.urlparse(url)

    # Keep valid HTTPS URLs unchanged.
    if parsed.scheme == "https":
        return url

    # Auto-upgrade HTTP links to HTTPS.
    if parsed.scheme == "http":
        return urllib.parse.urlunparse(parsed._replace(scheme="https"))

    # Reject non-web schemes (file:, ftp:, mailto:, etc.).
    raise WebToolError("Only https URLs are allowed.")


def _normalize_content_source_url(url: str) -> str:
    """
    Normalize source URLs for better content extraction.

    Example: GitHub blob URLs are mapped to raw.githubusercontent.com so
    code files can be fetched as plain text.
    """
    # Force HTTPS and normalize user input URL first.
    url = ensure_https_url(url)
    # Parse the URL so host/path can be inspected safely.
    try:
        parsed = urllib.parse.urlparse(url)
    # If parsing fails, return the normalized input unchanged.
    except Exception:
        return url

    # Normalize host to lowercase for reliable host matching.
    host = (parsed.netloc or "").lower()
    # Split path into non-empty segments.
    parts = [p for p in (parsed.path or "").split("/") if p]
    # Match GitHub blob links: /<owner>/<repo>/blob/<ref>/<path...>.
    if host.endswith("github.com") and len(parts) >= 5 and parts[2] == "blob":
        # Extract owner, repository, and git ref from path segments.
        owner, repo, ref = parts[0], parts[1], parts[3]
        # Rebuild remaining file path tail.
        tail = "/".join(parts[4:])
        # Convert to raw URL only when there is a file path tail.
        if tail:
            # Build raw.githubusercontent.com URL tuple:
            # (scheme, netloc, path, params, query, fragment).
            return urllib.parse.urlunparse(
                (
                    "https",
                    "raw.githubusercontent.com",
                    f"/{owner}/{repo}/{ref}/{tail}",
                    "",
                    "",
                    "",
                )
            )
    # Return original normalized URL when no rewrite rule applies.
    return url


class _HTTPSOnlyRedirectHandler(urllib.request.HTTPRedirectHandler):
    def redirect_request(self, req, fp, code, msg, headers, newurl):  # type: ignore[override]
        # Resolve relative redirects against the current URL.
        absolute = urllib.parse.urljoin(req.full_url, newurl)
        # Enforce HTTPS on each redirect hop.
        absolute = ensure_https_url(absolute)
        # Reuse default redirect logic with normalized URL.
        return super().redirect_request(req, fp, code, msg, headers, absolute)


# Shared opener so every HTTP call uses HTTPS-only redirect behavior.
_OPENER = urllib.request.build_opener(_HTTPSOnlyRedirectHandler())


def _http_get_bytes(url: str, *, timeout_s: float = 20.0, max_bytes: int = 2_000_000) -> tuple[str, bytes]:
    # Build request headers to reduce bot-blocking on some sites.
    req = urllib.request.Request(
        url,
        headers={
            "User-Agent": "Mozilla/5.0 (compatible; LocalAI-Agent/1.0)",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        },
    )

    try:
        # Stream response in bounded chunks to avoid unbounded memory growth.
        with _OPENER.open(req, timeout=timeout_s) as resp:
            content_type = resp.headers.get("Content-Type", "") or ""
            chunks: list[bytes] = []
            total = 0
            while True:
                # Cap each read and never exceed max_bytes overall.
                chunk = resp.read(min(64_000, max_bytes - total))
                if not chunk:
                    break
                chunks.append(chunk)
                total += len(chunk)
                # Stop reading once the safety cap is reached.
                if total >= max_bytes:
                    break
    except urllib.error.HTTPError as e:
        # Surface HTTP status code in a single consistent error type.
        raise WebToolError(f"HTTP {e.code} while fetching {url}") from e
    except urllib.error.URLError as e:
        # Surface DNS/connectivity/timeout issues consistently.
        raise WebToolError(f"Network error while fetching {url}: {e.reason}") from e

    # Join all chunks into one payload for downstream parsers.
    raw = b"".join(chunks)
    # Return lowercased content type so caller checks are case-insensitive.
    return content_type.lower(), raw


def _http_get(url: str, *, timeout_s: float = 20.0, max_bytes: int = 2_000_000) -> tuple[str, str]:
    # Fetch bytes first; decoding is handled separately for clarity.
    content_type, raw = _http_get_bytes(url, timeout_s=timeout_s, max_bytes=max_bytes)
    # Pick best-effort encoding from headers + HTML meta.
    encoding = _guess_encoding(content_type, raw)
    # Decode with replacement to avoid hard failures on mixed encodings.
    text = raw.decode(encoding, errors="replace")
    return content_type, text


def _guess_encoding(content_type: str, raw: bytes) -> str:
    # Start with HTTP header charset when present.
    ct = content_type or ""
    match = re.search(r"charset=([^\s;]+)", ct, re.IGNORECASE)
    if match:
        return match.group(1).strip("\"'").lower()
    # If header is missing, inspect the HTML head for charset declarations.
    head = raw[:50_000].decode("utf-8", errors="ignore")
    match = re.search(r"charset\s*=\s*([a-zA-Z0-9_\-]+)", head, re.IGNORECASE)
    if match:
        return match.group(1).lower()
    # Default fallback for modern web content.
    return "utf-8"


class _HTMLTextExtractor(HTMLParser):
    # Lightweight HTML -> text extractor for fetched pages.
    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        # Output token buffer; joined and normalized in `text()`.
        self._out: list[str] = []
        # Depth counter for ignored blocks like <script>/<style>.
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        # Normalize tag casing once for comparisons.
        tag = tag.lower()
        # Enter skip mode for non-content tags.
        if tag in {"script", "style", "noscript"}:
            self._skip_depth += 1
            return
        # Ignore all content while inside skipped tags.
        if self._skip_depth:
            return
        # Insert logical breaks at common block-level tags.
        if tag in {"p", "div", "br", "hr", "li", "tr", "h1", "h2", "h3", "h4", "h5", "h6", "blockquote"}:
            self._out.append("\n")

    def handle_endtag(self, tag: str) -> None:
        # Normalize tag casing once for comparisons.
        tag = tag.lower()
        # Leave skip mode when the matching ignored block closes.
        if tag in {"script", "style", "noscript"} and self._skip_depth:
            self._skip_depth -= 1
            return
        # Ignore end tags while still inside skipped blocks.
        if self._skip_depth:
            return
        # Add a separator when closing block elements.
        if tag in {"p", "div", "li", "tr", "blockquote"}:
            self._out.append("\n")

    def handle_data(self, data: str) -> None:
        # Ignore text captured inside skipped blocks.
        if self._skip_depth:
            return
        # Drop pure-whitespace fragments.
        s = (data or "").strip()
        if not s:
            return
        # Keep text and add a space separator between tokens.
        self._out.append(s)
        self._out.append(" ")

    def text(self) -> str:
        # Join parser fragments into one raw string.
        raw = "".join(self._out)
        # Normalize whitespace/newlines for readable plain text output.
        raw = re.sub(r"[ \t]+\n", "\n", raw)
        raw = re.sub(r"\n{3,}", "\n\n", raw)
        raw = re.sub(r"[ \t]{2,}", " ", raw)
        return raw.strip()


# Fetch and convert supported web content into readable plain text.
def fetch_url(url: str, *, timeout_s: float = 20.0, max_chars: int = 8_000) -> str:
    # Rewrite known source URLs (e.g., GitHub blob -> raw).
    url = _normalize_content_source_url(url)
    # Download bytes once; content-type driven handlers parse from memory.
    content_type, raw = _http_get_bytes(url, timeout_s=timeout_s, max_bytes=12_000_000)

    # Precompute lowercase variants for robust type checks.
    lowered = (content_type or "").lower()
    url_lower = url.lower()

    # Final extracted text payload.
    text = ""
    # PDF branch.
    if "application/pdf" in lowered or url_lower.endswith(".pdf"):
        try:
            from pypdf import PdfReader  # type: ignore

            reader = PdfReader(io.BytesIO(raw))
            parts = []
            # Cap page extraction to keep latency bounded on very long PDFs.
            for i, page in enumerate(reader.pages):
                if i >= 60:
                    break
                try:
                    # Extract text per page; keep empty placeholder on failures.
                    parts.append(page.extract_text() or "")
                except Exception:
                    parts.append("")
            # Merge non-empty pages into one normalized text blob.
            text = "\n\n".join([p.strip() for p in parts if p and p.strip()]).strip()
        except Exception as e:
            raise WebToolError(f"Failed to extract PDF text from {url}: {e}") from e
    # DOCX branch.
    elif (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in lowered
        or url_lower.endswith(".docx")
    ):
        try:
            from docx import Document  # type: ignore

            doc = Document(io.BytesIO(raw))
            # Keep only non-empty paragraphs.
            paras = [p.text for p in doc.paragraphs if (p.text or "").strip()]
            text = "\n".join(paras).strip()
        except Exception as e:
            raise WebToolError(f"Failed to extract DOCX text from {url}: {e}") from e
    # PPTX branch.
    elif (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation" in lowered
        or url_lower.endswith(".pptx")
    ):
        try:
            from pptx import Presentation  # type: ignore

            pres = Presentation(io.BytesIO(raw))
            out: list[str] = []
            # Cap slides for predictable runtime.
            for si, slide in enumerate(pres.slides):
                if si >= 80:
                    break
                # Keep slide markers to preserve structure.
                out.append(f"[Slide {si+1}]")
                for shape in slide.shapes:
                    # Many PPTX shapes expose a plain `text` attribute.
                    s = (getattr(shape, "text", "") or "").strip()
                    if s:
                        out.append(s)
                out.append("")
            text = "\n".join(out).strip()
        except Exception as e:
            raise WebToolError(f"Failed to extract PPTX text from {url}: {e}") from e
    # XLSX/XLSM branch.
    elif (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" in lowered
        or url_lower.endswith(".xlsx")
        or url_lower.endswith(".xlsm")
    ):
        try:
            import openpyxl  # type: ignore

            wb = openpyxl.load_workbook(io.BytesIO(raw), data_only=True, read_only=True)
            out: list[str] = []
            # Cap sheets and rows to avoid huge outputs.
            for wi, name in enumerate(wb.sheetnames):
                if wi >= 6:
                    break
                ws = wb[name]
                out.append(f"[Sheet] {name}")
                for ri, row in enumerate(ws.iter_rows(values_only=True)):
                    if ri >= 200:
                        break
                    # Serialize first columns as tab-separated text.
                    cells = ["" if v is None else str(v) for v in row[:30]]
                    if any(c.strip() for c in cells):
                        out.append("\t".join(cells).rstrip())
                out.append("")
            text = "\n".join(out).strip()
        except Exception as e:
            raise WebToolError(f"Failed to extract XLSX text from {url}: {e}") from e
    # Image branch (OCR.Space).
    elif lowered.startswith("image/") or any(
        url_lower.endswith(ext) for ext in (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff")
    ):
        try:
            # Use configured OCR language if provided.
            lang = os.getenv("OCR_SPACE_LANGUAGE") or os.getenv("OCR_SPACE_LANG")
            # Derive a stable filename for MIME/type inference upstream.
            path = urllib.parse.urlparse(url).path or ""
            fname = (path.rsplit("/", 1)[-1] or "image").strip() or "image"
            text = ocr_api.ocr_image_bytes(raw, filename=fname, language=lang)
        except Exception as e:
            raise WebToolError(f"Failed to OCR image from {url}: {e}") from e
    else:
        # Generic text/HTML branch.
        encoding = _guess_encoding(content_type, raw)
        text = raw.decode(encoding, errors="replace")
        # Convert HTML to plain text; keep non-HTML as-is.
        if "text/html" in lowered or "<html" in text[:500].lower():
            parser = _HTMLTextExtractor()
            parser.feed(text)
            text = parser.text()
        else:
            text = text.strip()

    # Apply output length cap for terminal readability.
    if len(text) > max_chars:
        text = text[:max_chars].rstrip() + "..."
    return text


def _decode_ddg_href(href: str) -> str:
    # Normalize empty/None input.
    href = (href or "").strip()
    if not href:
        return ""
    # Expand scheme-relative and relative DuckDuckGo links.
    if href.startswith("//"):
        href = "https:" + href
    elif href.startswith("/duckduckgo.com/"):
        href = "https://" + href.lstrip("/")
    if href.startswith("/"):
        href = "https://duckduckgo.com" + href
    try:
        # Parse redirect wrapper links like /l/?uddg=...
        parsed = urllib.parse.urlparse(href)
        path = (parsed.path or "").rstrip("/")
        last_seg = path.rsplit("/", 1)[-1] if path else ""
        if parsed.netloc.endswith("duckduckgo.com") and last_seg == "l":
            qs = urllib.parse.parse_qs(parsed.query)
            if "uddg" in qs and qs["uddg"]:
                # Return decoded destination URL.
                return urllib.parse.unquote(qs["uddg"][0])
    except Exception:
        # On parser errors, return original href unchanged.
        return href
    return href


class _DuckDuckGoHTMLResults(HTMLParser):
    # Extract search result links from duckduckgo.com/html response.
    def __init__(self, *, max_results: int) -> None:
        super().__init__(convert_charrefs=True)
        self.max_results = max_results
        self.results: list[dict[str, str]] = []
        self._capturing_title = False
        self._current: dict[str, str] | None = None
        self._title_parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        # Stop parsing once caller-requested result cap is reached.
        if len(self.results) >= self.max_results:
            return
        # Only anchor tags can produce results.
        if tag.lower() != "a":
            return
        # Normalize attributes for case-insensitive lookup.
        attr_map = {k.lower(): (v or "") for k, v in attrs}
        cls = attr_map.get("class", "")
        href = attr_map.get("href", "")
        if not href:
            return
        # Match DuckDuckGo "result title link" anchors.
        if "result__a" in cls or "result-link" in cls:
            url = _decode_ddg_href(href)
            self._current = {"title": "", "url": url, "snippet": ""}
            self._title_parts = []
            self._capturing_title = True

    def handle_endtag(self, tag: str) -> None:
        # Finalize title capture at end of anchor.
        if tag.lower() != "a":
            return
        if self._capturing_title and self._current is not None:
            title = " ".join([p.strip() for p in self._title_parts if p.strip()]).strip()
            title = re.sub(r"\s{2,}", " ", title)
            self._current["title"] = title
            if self._current.get("url") and self._current.get("title"):
                self.results.append(self._current)
            self._current = None
            self._title_parts = []
            self._capturing_title = False

    def handle_data(self, data: str) -> None:
        if not self._capturing_title:
            return
        if self._current is None:
            return
        s = (data or "").strip()
        if s:
            self._title_parts.append(s)


class _DuckDuckGoLiteResults(HTMLParser):
    # Extract results from lite.duckduckgo.com HTML (simpler layout).
    def __init__(self, *, max_results: int) -> None:
        super().__init__(convert_charrefs=True)
        self.max_results = max_results
        self.results: list[dict[str, str]] = []
        self._capturing_title = False
        self._current: dict[str, str] | None = None
        self._title_parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        # Stop parsing once caller-requested result cap is reached.
        if len(self.results) >= self.max_results:
            return
        # Only anchors can represent outgoing result links.
        if tag.lower() != "a":
            return
        attr_map = {k.lower(): (v or "") for k, v in attrs}
        href = attr_map.get("href", "")
        if not href:
            return

        # Decode redirect wrappers and keep only real web URLs.
        url = _decode_ddg_href(href)
        try:
            parsed = urllib.parse.urlparse(url)
        except Exception:
            return
        if parsed.scheme not in {"http", "https"}:
            return
        if not parsed.netloc:
            return
        if parsed.netloc.endswith("duckduckgo.com"):
            return

        # Start title capture for this candidate result.
        self._current = {"title": "", "url": url, "snippet": ""}
        self._title_parts = []
        self._capturing_title = True

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() != "a":
            return
        if self._capturing_title and self._current is not None:
            title = " ".join([p.strip() for p in self._title_parts if p.strip()]).strip()
            title = re.sub(r"\s{2,}", " ", title)
            self._current["title"] = title
            if self._current.get("url") and self._current.get("title"):
                self.results.append(self._current)
            self._current = None
            self._title_parts = []
            self._capturing_title = False

    def handle_data(self, data: str) -> None:
        if not self._capturing_title:
            return
        if self._current is None:
            return
        s = (data or "").strip()
        if s:
            self._title_parts.append(s)


def _duckduckgo_html_search(query: str, *, max_results: int) -> list[dict[str, str]]:
    # Encode query for URL transport.
    q = urllib.parse.quote_plus(query)
    url = f"https://duckduckgo.com/html/?q={q}"
    # Fetch and parse HTML result page.
    content_type, page = _http_get(url, timeout_s=20.0, max_bytes=1_500_000)
    if "text/html" not in content_type and "<html" not in page[:300].lower():
        raise WebToolError("DuckDuckGo did not return HTML.")
    parser = _DuckDuckGoHTMLResults(max_results=max_results)
    parser.feed(page)
    return parser.results


def _duckduckgo_lite_search(query: str, *, max_results: int) -> list[dict[str, str]]:
    # Encode query for URL transport.
    q = urllib.parse.quote_plus(query)
    url = f"https://lite.duckduckgo.com/lite/?q={q}"
    # Fetch and parse lite HTML result page.
    content_type, page = _http_get(url, timeout_s=20.0, max_bytes=1_500_000)
    if "text/html" not in content_type and "<html" not in page[:300].lower():
        raise WebToolError("DuckDuckGo Lite did not return HTML.")
    parser = _DuckDuckGoLiteResults(max_results=max_results)
    parser.feed(page)
    return parser.results


def _duckduckgo_instant_answer(query: str, *, max_results: int) -> list[dict[str, str]]:
    # Encode query for URL transport.
    q = urllib.parse.quote_plus(query)
    url = f"https://api.duckduckgo.com/?q={q}&format=json&no_redirect=1&no_html=1&skip_disambig=1"
    # Fetch JSON payload from the instant-answer endpoint.
    _, text = _http_get(url, timeout_s=20.0, max_bytes=2_000_000)
    data = json.loads(text)
    out: list[dict[str, str]] = []

    # Prefer abstract answer when available.
    abstract = (data.get("AbstractText") or "").strip()
    abstract_url = (data.get("AbstractURL") or "").strip()
    if abstract_url:
        out.append({"title": (data.get("Heading") or "DuckDuckGo Abstract").strip(), "url": abstract_url, "snippet": abstract})

    def add_topic(item: dict[str, Any]) -> None:
        nonlocal out
        # Extract link+text from one topic row.
        url2 = (item.get("FirstURL") or "").strip()
        txt = (item.get("Text") or "").strip()
        if url2 and txt:
            out.append({"title": txt, "url": url2, "snippet": txt})

    # Traverse related topics (flat or nested groups).
    topics = data.get("RelatedTopics") or []
    for item in topics:
        if len(out) >= max_results:
            break
        if isinstance(item, dict) and "Topics" in item and isinstance(item["Topics"], list):
            for sub in item["Topics"]:
                if len(out) >= max_results:
                    break
                if isinstance(sub, dict):
                    add_topic(sub)
        elif isinstance(item, dict):
            add_topic(item)

    return out[:max_results]


# Perform web search with fallback strategies and return normalized results.
def web_search(query: str, *, max_results: int = 5) -> list[dict[str, str]]:
    # Normalize the incoming query and avoid `None`.
    query = (query or "").strip()
    # Fail fast when query is empty after trimming.
    if not query:
        raise WebToolError("Query is empty.")
    # Keep result count in a safe bounded range.
    max_results = max(1, min(int(max_results), 10))

    def normalize_results(results: list[dict[str, str]]) -> list[dict[str, str]]:
        # Final cleaned list returned to the caller.
        out: list[dict[str, str]] = []
        # Track URLs already emitted to avoid duplicates.
        seen: set[str] = set()
        # Walk each raw search item from provider parsers.
        for item in results:
            # Read and trim the candidate URL field.
            url = (item.get("url") or "").strip()
            # Skip rows with no URL.
            if not url:
                continue
            # Decode DuckDuckGo redirect wrappers to real destination URLs.
            url = _decode_ddg_href(url)
            # Enforce HTTPS-only output.
            try:
                url = ensure_https_url(url)
            except WebToolError:
                continue
            # Ignore links that still point to DuckDuckGo pages.
            try:
                parsed = urllib.parse.urlparse(url)
                if parsed.netloc.endswith("duckduckgo.com"):
                    continue
            except Exception:
                # Skip malformed URLs we cannot parse reliably.
                continue
            # Remove duplicates while preserving first occurrence order.
            if url in seen:
                continue
            # Mark URL as seen.
            seen.add(url)
            # Append normalized record with guaranteed keys.
            out.append(
                {
                    "title": (item.get("title") or "").strip(),
                    "url": url,
                    "snippet": (item.get("snippet") or "").strip(),
                }
            )
            # Stop once we reached the requested cap.
            if len(out) >= max_results:
                break
        # Return cleaned and bounded results.
        return out

    # Storage for the current provider attempt.
    results: list[dict[str, str]] = []
    # First attempt: DuckDuckGo HTML endpoint.
    try:
        results = normalize_results(_duckduckgo_html_search(query, max_results=max_results))
    except Exception:
        # Swallow provider failure and continue to fallback.
        results = []

    # Return immediately when first provider succeeded.
    if results:
        return results[:max_results]

    # Second attempt: DuckDuckGo Lite endpoint.
    try:
        results = normalize_results(_duckduckgo_lite_search(query, max_results=max_results))
    except Exception:
        # Swallow provider failure and continue to fallback.
        results = []

    # Return immediately when second provider succeeded.
    if results:
        return results[:max_results]

    # Final attempt: DuckDuckGo Instant Answer JSON API.
    try:
        results = normalize_results(_duckduckgo_instant_answer(query, max_results=max_results))
        # Return even if it is empty; this is the last fallback.
        return results[:max_results]
    except Exception as e:
        # Surface a single normalized error for callers.
        raise WebToolError(f"Search failed: {e}") from e
